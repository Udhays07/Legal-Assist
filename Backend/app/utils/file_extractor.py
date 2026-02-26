"""
app/utils/file_extractor.py

Utility for extracting plain text from uploaded files of various formats.

Supported formats
-----------------
- Plain text : .txt, .md, .csv, .log, .xml, .html, .htm, .json
- PDF         : .pdf  (via pdfminer.six — no Poppler required)
- Word        : .docx (via python-docx)
- Excel       : .xlsx (via openpyxl)
- PowerPoint  : .pptx (via python-pptx)
- Unknown     : Attempted as UTF-8; raises UnsupportedFileTypeError on failure

Constants
---------
MIN_CONTENT_LENGTH : int
    Minimum number of non-whitespace characters required for a document
    to be accepted.  Change this value to adjust the global default.
"""

from __future__ import annotations

import io
import os
import re

import chardet

MIN_CONTENT_LENGTH = 100  # characters of cleaned text

# ---------------------------------------------------------------------------
# Exceptions
# ---------------------------------------------------------------------------

class UnsupportedFileTypeError(Exception):
    """Raised when the file extension is not recognised and raw decoding fails."""


class ContentTooShortError(Exception):
    """Raised when extracted text does not meet the minimum length threshold."""

    def __init__(self, extracted_len: int, minimum: int) -> None:
        self.extracted_len = extracted_len
        self.minimum = minimum
        super().__init__(
            f"Extracted content is too short ({extracted_len} chars). "
            f"Minimum required: {minimum} chars."
        )


class FileExtractionError(Exception):
    """Raised when a supported file type fails to be parsed (e.g. corrupt file)."""


# ---------------------------------------------------------------------------
# Format-specific extractors (internal)
# ---------------------------------------------------------------------------

def _extract_plain_text(file_bytes: bytes) -> str:
    """Decode plain-text bytes, auto-detecting encoding with chardet."""
    try:
        detected = chardet.detect(file_bytes)
        encoding = detected.get("encoding") or "utf-8"
        return file_bytes.decode(encoding, errors="replace")
    except Exception as e:
        raise FileExtractionError(f"Failed to decode text: {str(e)}")


def _extract_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF using pdfminer.six."""
    try:
        from pdfminer.high_level import extract_text_to_fp
        from pdfminer.layout import LAParams
    except ImportError:  # pragma: no cover
        raise FileExtractionError("PDF extraction library (pdfminer.six) is not installed.")

    try:
        output = io.StringIO()
        extract_text_to_fp(io.BytesIO(file_bytes), output, laparams=LAParams())
        return output.getvalue()
    except Exception as e:
        raise FileExtractionError(f"Failed to parse PDF: {str(e)}")


def _extract_docx(file_bytes: bytes) -> str:
    """Extract text from a .docx Word document using python-docx."""
    try:
        import docx
    except ImportError:  # pragma: no cover
        raise FileExtractionError("Word extraction library (python-docx) is not installed.")

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(para.text for para in doc.paragraphs)
    except Exception as e:
        # Often caught for non-zip or corrupt .docx files
        raise FileExtractionError(f"Failed to parse Word document: {str(e)}")


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract text from an .xlsx spreadsheet using openpyxl."""
    try:
        import openpyxl
    except ImportError:  # pragma: no cover
        raise FileExtractionError("Excel extraction library (openpyxl) is not installed.")

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    lines.append(row_text)
        return "\n".join(lines)
    except Exception as e:
        raise FileExtractionError(f"Failed to parse Excel spreadsheet: {str(e)}")


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract text from a .pptx presentation using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover
        raise FileExtractionError("PowerPoint extraction library (python-pptx) is not installed.")

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            lines.append(text)
        return "\n".join(lines)
    except Exception as e:
        raise FileExtractionError(f"Failed to parse PowerPoint presentation: {str(e)}")


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

_PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".log", ".xml",
    ".html", ".htm", ".json",
}

_EXTRACTOR_MAP = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def _clean_text(raw: str) -> str:
    """Collapse consecutive whitespace and strip leading/trailing space.
    
    This also removes null characters (\\x00) which are rejected by PostgreSQL.
    """
    cleaned = raw.replace("\x00", "")
    return re.sub(r"\s+", " ", cleaned).strip()


def extract_text(filename: str, file_bytes: bytes) -> str:
    """
    Extract and return cleaned text from *file_bytes*.

    Parameters
    ----------
    filename : str
        Original file name (used to determine the format by extension).
    file_bytes : bytes
        Raw binary content of the uploaded file.

    Returns
    -------
    str
        Cleaned text extracted from the file.

    Raises
    ------
    UnsupportedFileTypeError
        If the file extension is not supported and raw decoding also fails.
    ContentTooShortError
        If the cleaned text length is below ``MIN_CONTENT_LENGTH``.
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext in _PLAIN_TEXT_EXTENSIONS:
        raw = _extract_plain_text(file_bytes)
    elif ext in _EXTRACTOR_MAP:
        raw = _EXTRACTOR_MAP[ext](file_bytes)
    else:
        # Unknown extension — attempt plain-text decode as a last resort.
        try:
            raw = _extract_plain_text(file_bytes)
            # Sanity-check: the result should be mostly printable characters.
            printable_ratio = sum(c.isprintable() for c in raw) / max(len(raw), 1)
            if printable_ratio < 0.70:
                raise UnsupportedFileTypeError(
                    f"File type '{ext}' is not supported and appears to be binary."
                )
        except UnicodeDecodeError:
            raise UnsupportedFileTypeError(
                f"File type '{ext}' is not supported."
            )

    cleaned = _clean_text(raw)

    if len(cleaned) < MIN_CONTENT_LENGTH:
        raise ContentTooShortError(len(cleaned), MIN_CONTENT_LENGTH)

    return cleaned
